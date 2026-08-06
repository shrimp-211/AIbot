"""多模态感知器:图片/音频/视频/文档/代码的统一理解入口。

参照 AstrBot 的 perception 能力 + mainidea.txt 多模态输入层设计。
各类感知器可独立实例化,也支持组合使用(如"分析这张图表并翻译")。
依赖均惰性导入:未安装对应库时,该感知器返回明确的能力缺失提示,不阻塞启动。
"""
from __future__ import annotations

import asyncio
import re
from pathlib import Path
from typing import Any

from loguru import logger

# 通用 MIME 类型判断(扩展名 → 类别)
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}
_AUDIO_EXTS = {".mp3", ".wav", ".ogg", ".m4a", ".flac", ".aac", ".amr"}
_VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv"}
_DOC_EXTS = {".pdf", ".txt", ".md", ".docx", ".pptx", ".xlsx", ".epub", ".html"}
_CODE_EXTS = {".py", ".js", ".ts", ".java", ".c", ".cpp", ".go", ".rs", ".rb",
              ".php", ".sh", ".sql", ".html", ".css", ".json", ".yaml", ".yml", ".xml"}


def classify_media(path: str) -> str:
    """按扩展名判断媒体类别:image/audio/video/document/code/unknown。"""
    ext = Path(path).suffix.lower()
    if ext in _IMAGE_EXTS:
        return "image"
    if ext in _AUDIO_EXTS:
        return "audio"
    if ext in _VIDEO_EXTS:
        return "video"
    if ext in _DOC_EXTS:
        return "document"
    if ext in _CODE_EXTS:
        return "code"
    return "unknown"


class ImagePerceiver:
    """图片理解:OCR(多后端自动降级)+ 视觉问答。

    - OCR: pytesseract → PaddleOCR → 报错(适配未装某库的环境)
    - 视觉问答: 经 LLM vision API(由外部传入的 analyzer 回调完成);
      无视觉模型时自动降级为 OCR 提取画面文字(返回原文并标注)。
    """

    async def ocr(self, image_path: str, lang: str = "chi_sim+eng") -> str:
        """OCR 识别文字。多后端按可用性自动降级:pytesseract → PaddleOCR。"""
        text = await self._ocr_pytesseract(image_path, lang)
        if text:
            return text
        text = await self._ocr_paddle(image_path)
        if text:
            return text
        return ""

    @staticmethod
    async def _ocr_pytesseract(image_path: str, lang: str) -> str:
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            return ""
        try:

            def _run() -> str:
                img = Image.open(image_path)
                return (pytesseract.image_to_string(img, lang=lang) or "").strip()

            return await asyncio.to_thread(_run)
        except Exception as exc:  # noqa: BLE001
            logger.warning("pytesseract OCR 失败: {}", exc)
            return ""

    @staticmethod
    async def _ocr_paddle(image_path: str) -> str:
        try:
            from paddleocr import PaddleOCR
        except ImportError:
            return ""
        try:

            def _run() -> str:
                ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
                result = ocr.ocr(image_path, cls=True)
                lines: list[str] = []
                for page in result or []:
                    if not page:
                        continue
                    for item in page:
                        txt = (item or [None, None])[1]
                        if isinstance(txt, tuple) and txt:
                            lines.append(str(txt[0]))
                        elif txt:
                            lines.append(str(txt))
                return "\n".join(lines).strip()

            return await asyncio.to_thread(_run)
        except Exception as exc:  # noqa: BLE001
            logger.warning("PaddleOCR 失败: {}", exc)
            return ""

    async def describe(
        self,
        image_path: str,
        question: str = "请详细描述这张图片的内容。",
        analyzer=None,
    ) -> str:
        """视觉问答。analyzer 负责把图片路径+问题交给多模态 LLM。

        适配度保障:analyzer 缺失(无视觉模型)时,降级为 OCR 提取画面文字,
        保证"至少能读图",而非直接报错。
        """
        if analyzer is None:
            text = await self.ocr(image_path)
            if text:
                return f"(当前模型无图像能力,已用 OCR 提取画面文字)\n{text}"
            raise RuntimeError(
                "当前模型无图像能力,且本机未装 OCR 引擎"
                "(pytesseract / PaddleOCR 至少其一,或配置支持图像的 LLM)。"
            )
        return await analyzer(image_path, question)


class AudioPerceiver:
    """音频理解:语音转文字(STT)。"""

    def __init__(self, stt_provider=None):
        self.stt_provider = stt_provider

    async def transcribe(self, audio_path: str) -> str:
        if self.stt_provider is None:
            raise RuntimeError(
                "语音转文字需要配置 STT Provider(config.yaml 的 provider_stt 段,"
                "默认 Whisper)。"
            )
        try:
            return await self.stt_provider.transcribe(audio_path)
        except FileNotFoundError:
            raise
        except Exception as exc:  # noqa: BLE001
            logger.warning("语音转写失败: {}", exc)
            raise RuntimeError(f"语音转写失败: {exc}") from exc


class VideoPerceiver:
    """视频理解:抽取关键帧 → 逐帧描述(经 LLM)→ 汇总。

    opencv 用于帧提取,pytesseract 可选用于画面内文字。
    """

    async def extract_frames(self, video_path: str, max_frames: int = 4) -> list[str]:
        try:
            import cv2
        except ImportError as exc:
            raise RuntimeError(
                "视频理解需要 opencv-python(pip install opencv-python-headless)"
            ) from exc

        def _run() -> list[str]:
            cap = cv2.VideoCapture(video_path)
            if not cap.isOpened():
                raise RuntimeError(f"无法打开视频: {video_path}")
            total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            n = min(max_frames, max(1, total))
            frames: list[str] = []
            import os
            import tempfile

            tmpdir = tempfile.mkdtemp(prefix="qqvideo_")
            for i in range(n):
                pos = int(total * i / n) if total > 0 else i
                cap.set(cv2.CAP_PROP_POS_FRAMES, pos)
                ok, frame = cap.read()
                if not ok:
                    continue
                out = os.path.join(tmpdir, f"frame_{i}.jpg")
                cv2.imwrite(out, frame)
                frames.append(out)
            cap.release()
            return frames

        try:
            return await asyncio.to_thread(_run)
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(f"视频帧提取失败: {exc}") from exc

    async def summarize(self, video_path: str, analyzer=None) -> str:
        frames = await self.extract_frames(video_path)
        if not frames:
            return "无法从视频中提取到可用帧。"
        if analyzer is None:
            return f"已从视频提取 {len(frames)} 个关键帧(需多模态 LLM 描述画面)。"
        descriptions = []
        for i, f in enumerate(frames):
            try:
                descriptions.append(f"[帧{i+1}] {await analyzer(f, '描述这个视频画面内容')}")
            except Exception as exc:  # noqa: BLE001
                descriptions.append(f"[帧{i+1}] (分析失败: {exc})")
        return "\n".join(descriptions)


class DocumentPerceiver:
    """文档理解:PDF/Word/PPT/Excel/EPUB/纯文本 解析为可检索文本。"""

    async def parse(self, doc_path: str) -> dict:
        ext = Path(doc_path).suffix.lower()
        handler = {
            ".pdf": self._parse_pdf,
            ".docx": self._parse_docx,
            ".pptx": self._parse_pptx,
            ".xlsx": self._parse_xlsx,
            ".epub": self._parse_epub,
            ".html": self._parse_html,
        }.get(ext, self._parse_text)
        try:
            text = await handler(doc_path)
        except ImportError as exc:
            raise RuntimeError(
                f"解析 {ext} 需要额外依赖: {exc}(按需 pip install 后重试)"
            ) from exc
        return {
            "path": doc_path,
            "type": ext,
            "content": text,
            "char_count": len(text),
        }

    @staticmethod
    async def _parse_text(path: str) -> str:
        text = Path(path).read_text(encoding="utf-8", errors="replace")
        return text[:200_000]  # 限制单文档 20 万字符

    @staticmethod
    async def _parse_pdf(path: str) -> str:
        try:
            import fitz  # PyMuPDF
        except ImportError as exc:
            raise ImportError("pip install PyMuPDF") from exc

        def _run() -> str:
            doc = fitz.open(path)
            parts = []
            for i, page in enumerate(doc):
                if i >= 50:
                    break
                parts.append(page.get_text())
            doc.close()
            return "\n".join(parts)[:200_000]

        return await asyncio.to_thread(_run)

    @staticmethod
    async def _parse_docx(path: str) -> str:
        try:
            from docx import Document
        except ImportError as exc:
            raise ImportError("pip install python-docx") from exc

        def _run() -> str:
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)[:200_000]

        return await asyncio.to_thread(_run)

    @staticmethod
    async def _parse_pptx(path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError("pip install python-pptx") from exc

        def _run() -> str:
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides):
                if i >= 50:
                    break
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)[:200_000]

        return await asyncio.to_thread(_run)

    @staticmethod
    async def _parse_xlsx(path: str) -> str:
        try:
            import openpyxl
        except ImportError as exc:
            raise ImportError("pip install openpyxl") from exc

        def _run() -> str:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets[:10]:
                parts.append(f"[工作表 {ws.title}]")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join(str(c) for c in row if c is not None))
            wb.close()
            return "\n".join(parts)[:200_000]

        return await asyncio.to_thread(_run)

    @staticmethod
    async def _parse_epub(path: str) -> str:
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise ImportError("pip install EbookLib beautifulsoup4") from exc

        def _run() -> str:
            book = epub.read_epub(path)
            parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BeautifulSoup(item.get_content(), "html.parser")
                parts.append(soup.get_text("\n"))
            return "\n".join(parts)[:200_000]

        return await asyncio.to_thread(_run)

    @staticmethod
    async def _parse_html(path: str) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise ImportError("pip install beautifulsoup4") from exc

        def _run() -> str:
            html = Path(path).read_text(encoding="utf-8", errors="replace")
            soup = BeautifulSoup(html, "html.parser")
            return re.sub(r"\n{3,}", "\n\n", soup.get_text("\n")).strip()[:200_000]

        return await asyncio.to_thread(_run)


class CodePerceiver:
    """代码理解:语言识别 + 结构概览 + 函数签名提取(纯文本分析,无需 AST 库)。"""

    _SHEBANG_RE = re.compile(r"^#!.*\b(python|node|bash|sh|ruby|perl)\b", re.I)

    def analyze(self, code: str, filename: str = "") -> dict:
        lang = self._detect_language(code, filename)
        lines = code.splitlines()
        functions = self._extract_functions(lines, lang)
        return {
            "language": lang,
            "line_count": len(lines),
            "function_count": len(functions),
            "functions": functions[:50],
            "preview": "\n".join(lines[:40]),
        }

    def _detect_language(self, code: str, filename: str) -> str:
        if filename:
            ext = Path(filename).suffix.lower()
            mapping = {
                ".py": "python", ".js": "javascript", ".ts": "typescript",
                ".java": "java", ".c": "c", ".cpp": "cpp", ".go": "go",
                ".rs": "rust", ".rb": "ruby", ".php": "php", ".sh": "bash",
                ".sql": "sql", ".html": "html", ".css": "css",
                ".json": "json", ".yaml": "yaml", ".yml": "yaml", ".xml": "xml",
            }
            if ext in mapping:
                return mapping[ext]
        m = self._SHEBANG_RE.search(code)
        if m:
            return m.group(1).lower()
        return "unknown"

    @staticmethod
    def _extract_functions(lines: list[str], lang: str) -> list[str]:
        pats = {
            "python": re.compile(r"^(def|async def|class)\s+(\w+)"),
            "javascript": re.compile(r"^(function\s+(\w+)|(?:const|let|var)\s+(\w+)\s*=\s*(?:async\s*)?\(|class\s+(\w+))"),
            "typescript": re.compile(r"^(function\s+(\w+)|(?:const|let)\s+(\w+)\s*[:=]|class\s+(\w+)|interface\s+(\w+))"),
            "java": re.compile(r"^\s*(public|private|protected)?\s*(static\s+)?[\w<>\[\]]+\s+(\w+)\s*\("),
            "go": re.compile(r"^func\s+(\(\s*\w+\s+\*?\w+\s*\)\s*)?(\w+)\s*\("),
            "rust": re.compile(r"^fn\s+(\w+)"),
        }
        pat = pats.get(lang)
        if not pat:
            return []
        funcs = []
        for line in lines:
            m = pat.search(line)
            if m:
                name = next((g for g in m.groups() if g), "")
                funcs.append(f"{name} — {line.strip()[:80]}")
        return funcs


class PerceptionManager:
    """统一感知入口:按文件类别路由到对应感知器。"""

    def __init__(self, *, stt_provider=None, llm_analyzer=None):
        self.image = ImagePerceiver()
        self.audio = AudioPerceiver(stt_provider=stt_provider)
        self.video = VideoPerceiver()
        self.document = DocumentPerceiver()
        self.code = CodePerceiver()
        self._llm_analyzer = llm_analyzer

    def set_llm_analyzer(self, analyzer) -> None:
        """注入多模态 LLM 回调(image_path, question) -> str。"""
        self._llm_analyzer = analyzer

    async def analyze(self, file_path: str, question: str | None = None) -> dict:
        """按文件类型自动路由到最合适的感知能力。"""
        kind = classify_media(file_path)
        q = question or ""
        if kind == "image":
            text = await self.image.describe(file_path, q or "请描述这张图片内容", self._llm_analyzer)
            return {"kind": "image", "result": text}
        if kind == "audio":
            text = await self.audio.transcribe(file_path)
            return {"kind": "audio", "result": text}
        if kind == "video":
            text = await self.video.summarize(file_path, self._llm_analyzer)
            return {"kind": "video", "result": text}
        if kind == "document":
            doc = await self.document.parse(file_path)
            return {"kind": "document", "result": doc["content"][:2000], "meta": {k: v for k, v in doc.items() if k != "content"}}
        if kind == "code":
            code = Path(file_path).read_text(encoding="utf-8", errors="replace")
            return {"kind": "code", "result": self.code.analyze(code, file_path)}
        return {"kind": "unknown", "result": f"暂不支持分析该文件类型: {file_path}"}
